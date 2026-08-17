import io
import tempfile
from pathlib import Path
from unittest.mock import AsyncMock

import docx
import pptx
import pymupdf
import pytest

from orchestrator.parse import (
    DocumentParseError,
    parse_document,
    parse_document_with_ocr_fallback,
    render_document_context,
)


def _write_docx(path: Path, headings_and_body: list[tuple[str | None, str]]) -> None:
    document = docx.Document()
    for heading, body in headings_and_body:
        if heading is not None:
            document.add_heading(heading, level=1)
        if body:
            document.add_paragraph(body)
    document.save(str(path))


def _write_pptx(path: Path, slides: list[list[str]]) -> None:
    presentation = pptx.Presentation()
    for texts in slides:
        if texts:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = texts[0]
            if len(texts) > 1 and len(slide.placeholders) > 1:
                slide.placeholders[1].text_frame.text = "\n".join(texts[1:])
        else:
            presentation.slides.add_slide(presentation.slide_layouts[6])  # blank layout
    presentation.save(str(path))


def _write_pdf(path: Path, pages: list[str]) -> None:
    document = pymupdf.open()
    for text in pages:
        page = document.new_page()
        if text:
            page.insert_text((72, 72), text)
    document.save(str(path))
    document.close()


def _tiny_png_bytes() -> bytes:
    document = pymupdf.open()
    page = document.new_page(width=50, height=50)
    png_bytes = page.get_pixmap().tobytes("png")
    document.close()
    return png_bytes


def _write_image_only_docx(path: Path) -> None:
    document = docx.Document()
    document.add_picture(io.BytesIO(_tiny_png_bytes()))
    document.save(str(path))


def _write_pptx_image_only_slide(path: Path) -> None:
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(io.BytesIO(_tiny_png_bytes()), 0, 0)
    presentation.save(str(path))


class TestParseDocx:
    def test_headings_split_into_sections(self, tmp_path):
        path = tmp_path / "doc.docx"
        _write_docx(path, [("Executive Summary", "Summary text."), ("Findings", "Finding text.")])

        sections = parse_document(path)

        assert [s.section for s in sections] == ["Executive Summary", "Findings"]
        assert sections[0].text == "Summary text."
        assert sections[1].text == "Finding text."

    def test_no_headings_falls_back_to_single_document_section(self, tmp_path):
        path = tmp_path / "doc.docx"
        _write_docx(path, [(None, "Just a paragraph, no headings.")])

        sections = parse_document(path)

        assert len(sections) == 1
        assert sections[0].section == "Document"
        assert "Just a paragraph" in sections[0].text

    def test_empty_document_raises(self, tmp_path):
        path = tmp_path / "doc.docx"
        _write_docx(path, [])

        with pytest.raises(DocumentParseError, match="No readable text"):
            parse_document(path)

    def test_corrupted_file_raises_clean_error(self, tmp_path):
        path = tmp_path / "doc.docx"
        path.write_bytes(b"not a real docx file")

        with pytest.raises(DocumentParseError, match="Could not open this file as a Word document"):
            parse_document(path)

    def test_pages_are_none_for_docx(self, tmp_path):
        path = tmp_path / "doc.docx"
        _write_docx(path, [("H", "body")])

        sections = parse_document(path)

        assert sections[0].page is None


class TestParsePptx:
    def test_slides_become_sections_with_page_numbers(self, tmp_path):
        path = tmp_path / "deck.pptx"
        _write_pptx(path, [["Slide One", "body one"], ["Slide Two", "body two"]])

        sections = parse_document(path)

        assert len(sections) == 2
        assert sections[0].section == "Slide 1: Slide One"
        assert sections[0].page == 1
        assert sections[1].section == "Slide 2: Slide Two"
        assert sections[1].page == 2

    def test_slide_with_no_text_shapes_gets_placeholder_section(self, tmp_path):
        path = tmp_path / "deck.pptx"
        _write_pptx(path, [[], ["Slide Two", "body"]])

        sections = parse_document(path)

        assert len(sections) == 2
        assert sections[0].section == "Slide 1"
        assert "no text content" in sections[0].text
        assert sections[1].section == "Slide 2: Slide Two"

    def test_all_slides_blank_raises(self, tmp_path):
        path = tmp_path / "deck.pptx"
        _write_pptx(path, [[], []])

        # All-blank slides still produce placeholder sections (deck has structure,
        # even if no text) — this should NOT raise, unlike a truly empty docx.
        sections = parse_document(path)
        assert len(sections) == 2

    def test_corrupted_file_raises_clean_error(self, tmp_path):
        path = tmp_path / "deck.pptx"
        path.write_bytes(b"not a real pptx file")

        with pytest.raises(DocumentParseError, match="Could not open this file as a PowerPoint deck"):
            parse_document(path)


class TestParsePdf:
    def test_pages_become_sections(self, tmp_path):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["Page one text", "Page two text"])

        sections = parse_document(path)

        assert len(sections) == 2
        assert sections[0].section == "Page 1"
        assert sections[0].page == 1
        assert "Page one text" in sections[0].text

    def test_blank_pages_are_skipped(self, tmp_path):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["Real content here", ""])

        sections = parse_document(path)

        assert len(sections) == 1
        assert sections[0].section == "Page 1"

    def test_all_blank_pages_raises(self, tmp_path):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["", ""])

        with pytest.raises(DocumentParseError, match="No readable text"):
            parse_document(path)

    def test_corrupted_file_raises_clean_error(self, tmp_path):
        path = tmp_path / "doc.pdf"
        path.write_bytes(b"not a real pdf file")

        with pytest.raises(DocumentParseError, match="Could not open this file as a PDF"):
            parse_document(path)


class TestPdfOcrFallback:
    async def test_falls_back_to_vision_when_no_text_extracted(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["", ""])
        mock_transcribe = AsyncMock(return_value="Transcribed text.")
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 2
        assert {s.page for s in sections} == {1, 2}
        assert all(s.text == "Transcribed text." for s in sections)
        assert all(s.section == f"Page {s.page}" for s in sections)
        assert mock_transcribe.await_count == 2

    async def test_page_with_empty_transcription_is_skipped(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["", ""])
        # AsyncMock with a plain-list side_effect resolves synchronously (no real
        # suspension point), and gather() schedules tasks in creation order with
        # nothing to interleave them -- so this deterministically maps "" to page 1
        # and the real text to page 2, not by luck.
        mock_transcribe = AsyncMock(side_effect=["", "Transcribed second page."])
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert sections[0].page == 2
        assert sections[0].text == "Transcribed second page."

    async def test_raises_when_all_transcriptions_are_empty(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["", ""])
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", AsyncMock(return_value=""))

        with pytest.raises(DocumentParseError, match="No readable text"):
            await parse_document_with_ocr_fallback(path, client=AsyncMock())

    async def test_unsupported_extension_reraises_without_calling_client(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.txt"
        path.write_text("hello")
        mock_transcribe = AsyncMock()
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", mock_transcribe)

        with pytest.raises(DocumentParseError, match="Unsupported file type"):
            await parse_document_with_ocr_fallback(path, client=AsyncMock())
        mock_transcribe.assert_not_awaited()


class TestDocxOcrFallback:
    async def test_falls_back_to_vision_for_image_only_docx(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.docx"
        _write_image_only_docx(path)
        mock_transcribe = AsyncMock(return_value="Transcribed docx image.")
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert sections[0].section == "Image 1"
        assert sections[0].page is None
        assert sections[0].text == "Transcribed docx image."
        mock_transcribe.assert_awaited_once()

    async def test_raises_when_docx_image_transcription_is_empty(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.docx"
        _write_image_only_docx(path)
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", AsyncMock(return_value=""))

        with pytest.raises(DocumentParseError, match="No readable text"):
            await parse_document_with_ocr_fallback(path, client=AsyncMock())

    async def test_completely_empty_docx_raises_without_calling_client(self, tmp_path, monkeypatch):
        # No text AND no images at all -- distinct from the "image with empty
        # transcription" case above: here transcribe_page_image should never even
        # be reached, since there's nothing to send it. This also confirms the
        # related_parts fix: a blank docx still contains docProps/thumbnail.jpeg
        # (an image/* part) at the package level, which must NOT be picked up.
        path = tmp_path / "doc.docx"
        _write_docx(path, [])
        mock_transcribe = AsyncMock()
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", mock_transcribe)

        with pytest.raises(DocumentParseError, match="No readable text"):
            await parse_document_with_ocr_fallback(path, client=AsyncMock())
        mock_transcribe.assert_not_awaited()


class TestPptxOcrFallback:
    async def test_enriches_image_only_slide_via_vision(self, tmp_path, monkeypatch):
        path = tmp_path / "deck.pptx"
        _write_pptx_image_only_slide(path)
        mock_transcribe = AsyncMock(return_value="Transcribed slide image.")
        monkeypatch.setattr("orchestrator.parse.transcribe_page_images", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert sections[0].section == "Slide 1"
        assert sections[0].page == 1
        assert sections[0].text == "Transcribed slide image."
        mock_transcribe.assert_awaited_once()

    async def test_real_text_slides_are_never_sent_to_vision(self, tmp_path, monkeypatch):
        path = tmp_path / "deck.pptx"
        _write_pptx(path, [["Slide One", "body"]])
        mock_transcribe = AsyncMock()
        monkeypatch.setattr("orchestrator.parse.transcribe_page_images", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert sections[0].section == "Slide 1: Slide One"
        mock_transcribe.assert_not_awaited()

    async def test_blank_slide_with_no_picture_keeps_placeholder(self, tmp_path, monkeypatch):
        path = tmp_path / "deck.pptx"
        _write_pptx(path, [[]])  # blank layout -- no text, no picture either
        mock_transcribe = AsyncMock()
        monkeypatch.setattr("orchestrator.parse.transcribe_page_images", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert "no text content" in sections[0].text
        mock_transcribe.assert_not_awaited()

    async def test_slide_transcription_coming_back_empty_keeps_placeholder(self, tmp_path, monkeypatch):
        path = tmp_path / "deck.pptx"
        _write_pptx_image_only_slide(path)
        monkeypatch.setattr("orchestrator.parse.transcribe_page_images", AsyncMock(return_value=""))

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert "no text content" in sections[0].text

    async def test_successful_normal_parse_never_touches_the_client(self, tmp_path, monkeypatch):
        path = tmp_path / "doc.pdf"
        _write_pdf(path, ["Real text content"])
        mock_transcribe = AsyncMock()
        monkeypatch.setattr("orchestrator.parse.transcribe_page_image", mock_transcribe)

        sections = await parse_document_with_ocr_fallback(path, client=AsyncMock())

        assert len(sections) == 1
        assert sections[0].text == "Real text content"
        mock_transcribe.assert_not_awaited()


class TestParseDocument:
    def test_unsupported_extension_raises(self, tmp_path):
        path = tmp_path / "doc.txt"
        path.write_text("hello")

        with pytest.raises(DocumentParseError, match="Unsupported file type"):
            parse_document(path)


class TestRenderDocumentContext:
    def test_includes_all_inputs(self, tmp_path):
        path = tmp_path / "doc.docx"
        _write_docx(path, [("Heading", "Body text.")])
        sections = parse_document(path)

        context = render_document_context(sections, "advisory", "checklist: yaml", "style: yaml")

        assert "engagement_type: advisory" in context
        assert "checklist: yaml" in context
        assert "style: yaml" in context
        assert "=== Heading ===" in context
        assert "Body text." in context
