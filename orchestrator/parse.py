import asyncio
import base64
from dataclasses import dataclass
from pathlib import Path

import docx
import pptx
import pymupdf
from anthropic import AsyncAnthropicBedrock
from pptx.enum.shapes import MSO_SHAPE_TYPE

from agents.schema import transcribe_page_image, transcribe_page_images

# The literal text parse_pptx() gives a slide with no text shapes -- shared with
# ocr_image_only_slides() below, which looks for this exact marker to decide
# which already-parsed sections are candidates for vision OCR.
PPTX_BLANK_SLIDE_TEXT = "(no text content on this slide)"

# A US Letter page at this DPI renders to ~1275x1650px, already close to
# Anthropic's own recommended long-side (~1568px) for vision cost/quality
# balance -- no extra resize step needed.
OCR_RENDER_DPI = 150
# Caps simultaneous Bedrock requests when OCR-ing a fully-scanned deliverable
# -- an 80-page document would otherwise fire 80 requests at once.
OCR_CONCURRENCY_LIMIT = 5


class DocumentParseError(ValueError):
    """Raised when a document can't be parsed or has no extractable text.

    Always safe to show the message directly to an end user (e.g. in the
    web UI's error banner) — it never leaks a library traceback.
    """


@dataclass
class Section:
    section: str
    page: int | None
    text: str


def _open_docx(path: Path) -> docx.Document:
    try:
        return docx.Document(str(path))
    except Exception as e:
        raise DocumentParseError(f"Could not open this file as a Word document: {e}") from e


def parse_docx(path: Path) -> list[Section]:
    document = _open_docx(path)

    sections: list[Section] = []
    current_heading = "Document"
    buffer: list[str] = []

    def flush() -> None:
        if buffer:
            sections.append(Section(section=current_heading, page=None, text="\n".join(buffer)))
            buffer.clear()

    for para in document.paragraphs:
        if not para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            flush()
            current_heading = para.text.strip()
        else:
            buffer.append(para.text)
    flush()
    return sections


def parse_pptx(path: Path) -> list[Section]:
    try:
        presentation = pptx.Presentation(str(path))
    except Exception as e:
        raise DocumentParseError(f"Could not open this file as a PowerPoint deck: {e}") from e

    sections: list[Section] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            title = texts[0]
            body = "\n".join(texts)
            sections.append(Section(section=f"Slide {index}: {title}", page=index, text=body))
        else:
            sections.append(Section(section=f"Slide {index}", page=index, text=PPTX_BLANK_SLIDE_TEXT))
    return sections


def _open_pdf(path: Path) -> pymupdf.Document:
    try:
        return pymupdf.open(str(path))
    except Exception as e:
        raise DocumentParseError(f"Could not open this file as a PDF: {e}") from e


def parse_pdf(path: Path) -> list[Section]:
    document = _open_pdf(path)

    sections: list[Section] = []
    with document:
        for page_number, page in enumerate(document, start=1):
            text = page.get_text().strip()
            if text:
                sections.append(Section(section=f"Page {page_number}", page=page_number, text=text))
    return sections


async def ocr_scanned_pdf(path: Path, client: AsyncAnthropicBedrock) -> list[Section]:
    document = _open_pdf(path)
    semaphore = asyncio.Semaphore(OCR_CONCURRENCY_LIMIT)

    async def transcribe(page_number: int, page: pymupdf.Page) -> Section | None:
        image_b64 = base64.b64encode(page.get_pixmap(dpi=OCR_RENDER_DPI).tobytes("png")).decode()
        async with semaphore:
            text = await transcribe_page_image(client, image_b64)
        return Section(section=f"Page {page_number}", page=page_number, text=text) if text else None

    with document:
        results = await asyncio.gather(*(transcribe(n, page) for n, page in enumerate(document, start=1)))

    sections = [s for s in results if s is not None]
    if not sections:
        raise DocumentParseError(
            "No readable text was found in this document, even after attempting OCR via "
            "Claude vision. It may be blank, corrupted, or too low-resolution to read."
        )
    return sections


async def ocr_scanned_docx(path: Path, client: AsyncAnthropicBedrock) -> list[Section]:
    """For a .docx with zero extractable text (e.g. a scanned page pasted in as an
    image) -- transcribes every embedded image in the package, since docx has no
    page concept to key sections off the way pptx (slides) and pdf (pages) do."""
    document = _open_docx(path)
    # related_parts (not package.iter_parts()) -- scoped to what word/document.xml itself
    # references, so it excludes package-level assets like docProps/thumbnail.jpeg, which
    # is an image/* part present in every docx, blank or not.
    image_parts = [part for part in document.part.related_parts.values() if part.content_type.startswith("image/")]
    semaphore = asyncio.Semaphore(OCR_CONCURRENCY_LIMIT)

    async def transcribe(index: int, part) -> Section | None:
        image_b64 = base64.b64encode(part.blob).decode()
        async with semaphore:
            text = await transcribe_page_image(client, image_b64, media_type=part.content_type)
        return Section(section=f"Image {index}", page=None, text=text) if text else None

    results = await asyncio.gather(*(transcribe(n, part) for n, part in enumerate(image_parts, start=1)))

    sections = [s for s in results if s is not None]
    if not sections:
        raise DocumentParseError(
            "No readable text was found in this document, even after attempting OCR via "
            "Claude vision on its embedded images. It may be blank, corrupted, or too "
            "low-resolution to read."
        )
    return sections


async def ocr_image_only_slides(path: Path, sections: list[Section], client: AsyncAnthropicBedrock) -> list[Section]:
    """Best-effort enrichment for an already-successful pptx parse: parse_pptx() never
    raises (a slide with no text still gets a PPTX_BLANK_SLIDE_TEXT placeholder, since a
    deck having *some* image-only slides is normal, not a failure) -- this replaces just
    those placeholder sections with vision-transcribed text where the slide actually has
    a picture to transcribe, leaving every other section (and slides with genuinely no
    picture either) untouched."""
    try:
        presentation = pptx.Presentation(str(path))
    except Exception:
        return sections  # parse_document() already opened this file once successfully

    semaphore = asyncio.Semaphore(OCR_CONCURRENCY_LIMIT)

    async def enrich(section: Section) -> Section:
        if section.text != PPTX_BLANK_SLIDE_TEXT or section.page is None:
            return section
        slide = presentation.slides[section.page - 1]
        images = [
            (base64.b64encode(shape.image.blob).decode(), shape.image.content_type)
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        if not images:
            return section
        async with semaphore:
            text = await transcribe_page_images(client, images)
        return Section(section=section.section, page=section.page, text=text) if text else section

    return list(await asyncio.gather(*(enrich(s) for s in sections)))


def parse_document(path: Path) -> list[Section]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        sections = parse_docx(path)
    elif suffix == ".pptx":
        sections = parse_pptx(path)
    elif suffix == ".pdf":
        sections = parse_pdf(path)
    else:
        raise DocumentParseError(f"Unsupported file type {suffix!r} — expected .docx, .pptx, or .pdf")

    if not sections:
        raise DocumentParseError(
            "No readable text was found in this document. If it's a scanned PDF or an "
            "image-based deck, this pipeline can't extract text from it (OCR isn't supported)."
        )
    return sections


async def parse_document_with_ocr_fallback(path: Path, client: AsyncAnthropicBedrock) -> list[Section]:
    """Like parse_document(), but reaches for Claude-vision OCR instead of giving up on a
    scanned/image-only document -- see ocr_scanned_pdf(), ocr_scanned_docx(), and
    ocr_image_only_slides() for how each format's trigger condition differs (a pdf/docx
    with zero text raises and is retried; a pptx never raises in the first place, since a
    slide with no text already gets a placeholder section today, so it's enriched instead)."""
    suffix = path.suffix.lower()
    try:
        sections = parse_document(path)
    except DocumentParseError:
        if suffix == ".pdf":
            return await ocr_scanned_pdf(path, client)
        if suffix == ".docx":
            return await ocr_scanned_docx(path, client)
        raise

    if suffix == ".pptx":
        return await ocr_image_only_slides(path, sections, client)
    return sections


def render_document_context(sections: list[Section], engagement_type: str, checklist_yaml: str, style_rules_yaml: str) -> str:
    section_blocks = "\n\n".join(f"=== {s.section} ===\n{s.text}" for s in sections)
    return (
        f"engagement_type: {engagement_type}\n\n"
        f"--- checklist config ---\n{checklist_yaml}\n\n"
        f"--- style rules config ---\n{style_rules_yaml}\n\n"
        f"--- document sections ---\n{section_blocks}\n"
    )
